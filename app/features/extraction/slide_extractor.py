# backend/app/features/extraction/slide_extractor.py
import os
from pptx import Presentation
import pdfplumber

def extract_text_from_pptx(file_path: str) -> list:
    """Extracts structured text slide-by-slide from PPTX, handling titles, lists, and tables."""
    prs = Presentation(file_path)
    structured_slides = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_content = []
        slide_title = f"Slide {i}"

        # 1. Try to extract the official slide title shape
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide_title = slide.shapes.title.text.strip() or f"Slide {i}"

        # 2. Iterate through all shapes structurally
        for shape in slide.shapes:
            # Skip the title shape since we already captured it
            if shape == slide.shapes.title:
                continue

            # Process Text Boxes and Bullet Points
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    
                    # Indent nested bullets so the AI understands information hierarchy
                    indent = "  " * paragraph.level if paragraph.level > 0 else ""
                    bullet = "• " if paragraph.level >= 0 else ""
                    slide_content.append(f"{indent}{bullet}{text}")

            # Process Tables row-by-row so data rows don't get scrambled
            elif shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_cells:
                        table_data.append(" | ".join(row_cells))
                if table_data:
                    slide_content.append("\n[Table Data]:\n" + "\n".join(table_data))

        # Combine everything for this slide
        full_text = f"=== TITLE: {slide_title} ===\n" + "\n".join(slide_content)
        structured_slides.append({"slide": i, "text": full_text.strip()})

    return structured_slides


def extract_text_from_pdf(file_path: str) -> list:
    """Extracts text and layout-sensitive tables slide-by-slide from a PDF."""
    structured_slides = []

    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            # Extract regular page text
            text_content = page.extract_text(layout=False) or ""
            
            # Extract tables cleanly and format them as readable text strings
            tables = page.extract_tables()
            table_strings = []
            for table in tables:
                for row in table:
                    # Filter out None values and join cells
                    cleaned_row = [str(cell).strip() for cell in row if cell is not None]
                    if cleaned_row:
                        table_strings.append(" | ".join(cleaned_row))
            
            # Combine page text with formatted tables
            full_text = f"=== Slide {i} ===\n{text_content}"
            if table_strings:
                full_text += "\n\n[Table Data]:\n" + "\n".join(table_strings)

            structured_slides.append({"slide": i, "text": full_text.strip()})

    return structured_slides


def process_slide_file(file_path: str) -> list:
    """Master controller that auto-detects file extension and routes to the correct parser."""
    _, ext = os.path.splitext(file_path.lower())
    
    if ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.pdf':
        return extract_text_from_pdf(file_path)
    else:
        raise ValueError(f"Unsupported slide format: {ext}")